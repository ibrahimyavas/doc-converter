"""Document conversion via LibreOffice (headless), with two format-pair
classes routed around it instead:

- JPG: LibreOffice's own image import/export is unreliable for
  round-tripping quality, whereas Poppler (PDF -> image) and Pillow
  (image -> PDF) are purpose-built and predictable.
- PDF as the *source* for a text format (docx, txt, ...): LibreOffice
  opens a PDF through its Draw import filter (a flat visual
  reproduction), which has no export path to Writer formats at all —
  `soffice --convert-to docx some.pdf` fails outright with "no export
  filter", regardless of the PDF's actual content. Real text extraction
  needs `pdftotext` (Poppler) first; LibreOffice then converts *that*
  plain text into the target format, since a .txt file is something
  Writer actually opens natively.

Every function here takes a `workdir` and writes into it rather than
calling tempfile.mkdtemp() itself — a conversion can chain through
several intermediate files (jpg -> pdf -> txt -> pptx), and giving them
all one shared directory means the caller can delete the whole thing in
one shot once the response is sent, instead of every intermediate
leaking its own orphaned temp directory forever.
"""

import os
import shutil
import uuid
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from .tools import ConversionError, find_tool, run

# Formats LibreOffice converts between directly via --convert-to.
LIBREOFFICE_EXTS = {"pdf", "docx", "txt", "pptx", "xlsx", "html", "rtf", "odt", "epub"}

# LibreOffice is really three applications wearing one CLI: Writer, Calc,
# and Impress each read/write their own format family and *cannot*
# --convert-to a format outside it directly (Writer has no xlsx export,
# Calc has no docx export, etc. — there's no bridging filter, not just a
# missing extension mapping). Every one of them can export to PDF though,
# so that's the universal hop: same-family pairs convert directly,
# cross-family pairs go through PDF (native export from the source side,
# then _pdf_to_document's text-extraction/rebuild on the target side).
_WRITER_EXTS = {"docx", "txt", "html", "rtf", "odt", "epub"}
_CALC_EXTS = {"xlsx"}
_IMPRESS_EXTS = {"pptx"}


def _family(ext: str) -> str | None:
    if ext in _WRITER_EXTS:
        return "writer"
    if ext in _CALC_EXTS:
        return "calc"
    if ext in _IMPRESS_EXTS:
        return "impress"
    return None


def _soffice():
    return find_tool("SOFFICE_PATH", "soffice")


def _pdftoppm():
    return find_tool("PDFTOPPM_PATH", "pdftoppm")


def _pdftotext():
    return find_tool("PDFTOTEXT_PATH", "pdftotext")


def _new_dir(workdir: str, prefix: str) -> str:
    path = os.path.join(workdir, f"{prefix}_{uuid.uuid4().hex}")
    os.makedirs(path, exist_ok=True)
    return path


def _convert_via_libreoffice(input_path: str, to_ext: str, workdir: str) -> str:
    """Runs `soffice --convert-to` and returns the output file's path.

    Each call gets its own throwaway user-profile dir (via
    -env:UserInstallation) — LibreOffice locks its profile, so concurrent
    requests sharing one would collide with "another instance is running".
    """
    soffice = _soffice()
    outdir = _new_dir(workdir, "lo_out")
    profile_dir = _new_dir(workdir, "lo_profile")
    profile_uri = "file:///" + profile_dir.replace("\\", "/")

    run(
        [
            soffice,
            "--headless",
            "--norestore",
            f"-env:UserInstallation={profile_uri}",
            "--convert-to",
            to_ext,
            "--outdir",
            outdir,
            input_path,
        ],
        timeout=180,
    )

    produced = list(Path(outdir).glob(f"*.{to_ext}"))
    if not produced:
        raise ConversionError(f"LibreOffice didn't produce a .{to_ext} file — this pair may not be supported.")
    return str(produced[0])


def _pdf_to_jpg(pdf_path: str, workdir: str) -> str:
    pdftoppm = _pdftoppm()
    out_dir = _new_dir(workdir, "pdftoppm")
    out_prefix = os.path.join(out_dir, "page")
    # -f/-l 1: first page only — a multi-page PDF collapses to one image,
    # matching how every other single-file-out conversion here behaves.
    run([pdftoppm, "-jpeg", "-f", "1", "-l", "1", pdf_path, out_prefix], timeout=60)
    produced = sorted(Path(out_dir).glob("page*.jpg"))
    if not produced:
        raise ConversionError("pdftoppm didn't produce an image — the PDF may be empty or corrupt.")
    return str(produced[0])


def _jpg_to_pdf(jpg_path: str, workdir: str) -> str:
    out_path = os.path.join(workdir, f"{uuid.uuid4().hex}.pdf")
    try:
        with Image.open(jpg_path) as img:
            img.convert("RGB").save(out_path, "PDF")
    except Exception as exc:
        raise ConversionError(f"Could not read that image: {exc}") from exc
    return out_path


def _pdf_to_text(pdf_path: str, workdir: str) -> str:
    pdftotext = _pdftotext()
    out_path = os.path.join(workdir, f"{uuid.uuid4().hex}.txt")
    run([pdftotext, "-layout", pdf_path, out_path], timeout=60)
    if not os.path.isfile(out_path):
        raise ConversionError("pdftotext didn't produce any output.")
    return out_path


def _text_to_csv(txt_path: str, workdir: str) -> str:
    # LibreOffice picks its import filter off the file extension, not
    # content — Calc only recognizes plain text as data if it's *named*
    # .csv (a .txt with identical content opens in Writer instead, which
    # is a dead end for an xlsx export). One "cell" per line is a
    # reasonable, honest shape for unstructured extracted text.
    out_path = os.path.join(workdir, f"{uuid.uuid4().hex}.csv")
    shutil.copyfile(txt_path, out_path)
    return out_path


def _text_to_pptx(txt_path: str, workdir: str) -> str:
    # Writer (which opens .txt) has no export path to Impress either —
    # different application, no bridging filter. Building the deck
    # directly is simpler and more honest than hunting for one.
    with open(txt_path, encoding="utf-8", errors="replace") as f:
        text = f.read().strip()

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    chunk_size = 900  # characters per slide — keeps a textbox readable
    chunks = [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]

    for chunk in chunks:
        slide = prs.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = chunk
        tf.paragraphs[0].font.size = Pt(18)

    out_path = os.path.join(workdir, f"{uuid.uuid4().hex}.pptx")
    prs.save(out_path)
    return out_path


def _pdf_to_document(pdf_path: str, to_ext: str, workdir: str) -> str:
    """PDF -> any other document format. Always goes through pdftotext
    first (see module docstring). Where the target belongs to a
    different LibreOffice application than Writer (Calc for xlsx,
    Impress for pptx — Writer can't export directly to either), the
    extracted text is routed to a dedicated builder instead of
    LibreOffice's generic --convert-to.
    """
    txt_path = _pdf_to_text(pdf_path, workdir)
    if to_ext == "txt":
        return txt_path
    if to_ext == "xlsx":
        return _convert_via_libreoffice(_text_to_csv(txt_path, workdir), "xlsx", workdir)
    if to_ext == "pptx":
        return _text_to_pptx(txt_path, workdir)
    return _convert_via_libreoffice(txt_path, to_ext, workdir)


def convert_document(input_path: str, from_ext: str, to_ext: str, workdir: str) -> str:
    """Returns the path to the converted file, somewhere inside `workdir`
    — the caller owns workdir and should remove it (as a whole, e.g.
    shutil.rmtree) once the file's been used.
    """
    if from_ext == "jpg" and to_ext == "pdf":
        return _jpg_to_pdf(input_path, workdir)
    if from_ext == "pdf" and to_ext == "jpg":
        return _pdf_to_jpg(input_path, workdir)

    if from_ext == "jpg":
        # jpg -> (docx/txt/pptx/...): via an intermediate PDF, then real
        # text extraction. There's no actual text in a photo, so expect a
        # near-empty result — that's the honest answer without OCR, not
        # a bug.
        pdf_path = _jpg_to_pdf(input_path, workdir)
        return _pdf_to_document(pdf_path, to_ext, workdir)

    if to_ext == "jpg":
        # (docx/txt/pptx/...) -> jpg: source -> pdf -> jpg. This
        # direction is LibreOffice's well-supported native PDF *export*
        # — unlike PDF as a source, which is the problem this module
        # works around elsewhere.
        pdf_path = _convert_via_libreoffice(input_path, "pdf", workdir)
        return _pdf_to_jpg(pdf_path, workdir)

    if from_ext == "pdf":
        return _pdf_to_document(input_path, to_ext, workdir)

    if to_ext not in LIBREOFFICE_EXTS:
        raise ConversionError(f"Don't know how to produce .{to_ext} documents.")

    if _family(from_ext) == _family(to_ext):
        # Same application (e.g. docx -> html, both Writer) — LibreOffice
        # exports between these directly, no detour needed.
        return _convert_via_libreoffice(input_path, to_ext, workdir)

    # Cross-application (e.g. xlsx -> docx, pptx -> txt, docx -> pptx):
    # hop through PDF, which every LO application exports to natively,
    # then hand off to the same PDF-as-source logic used above.
    pdf_path = _convert_via_libreoffice(input_path, "pdf", workdir)
    return _pdf_to_document(pdf_path, to_ext, workdir)
